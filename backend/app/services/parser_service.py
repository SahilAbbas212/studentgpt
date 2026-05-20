import fitz
from docx import Document
from pptx import Presentation

def parse_pdf(path):
    text = ""

    doc = fitz.open(path)

    for page in doc:
        text += page.get_text()

    return text

def parse_docx(path):
    doc = Document(path)

    return "\n".join(
        [p.text for p in doc.paragraphs]
    )

def parse_pptx(path):
    prs = Presentation(path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text